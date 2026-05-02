import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def apply_theme(slide, theme_name):
    themes = {
        "light": {"bg": "#ffffff", "text": "#1d1d1f", "accent": "#034EA2"},
        "dark": {"bg": "#1d1d1f", "text": "#ffffff", "accent": "#2997ff"},
        "parchment": {"bg": "#f5f5f7", "text": "#1d1d1f", "accent": "#034EA2"}
    }
    theme = themes.get(theme_name, themes["light"])
    background = slide.background
    fill = background.fill
    fill.solid()
    r, g, b = hex_to_rgb(theme["bg"])
    fill.fore_color.rgb = RGBColor(r, g, b)
    return theme

def add_title(slide, text, theme):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.2))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Apple SD Gothic Neo'
    p.font.bold = True
    p.font.size = Pt(44)
    r, g, b = hex_to_rgb(theme["text"])
    p.font.color.rgb = RGBColor(r, g, b)

def format_text_in_paragraph(p, text):
    # PPT에서는 HTML태그를 바로 못쓰므로, 기호를 감지하여 색상을 변경하는 방식을 적용 가능하나,
    # 여기서는 심플하게 전체 텍스트를 넣습니다. (고급 렌더링은 향후 과제)
    p.text = text

def export_high_fidelity_ppt(json_data_path, output_pptx="presentation.pptx"):
    if not os.path.exists(json_data_path): return
    with open(json_data_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for item in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = apply_theme(slide, item.get('theme', 'light'))
        add_title(slide, item.get('title', ''), theme)
        stype = item.get('type')
        
        if stype == 'kpi':
            kpis = item.get('kpis', [])
            num_kpis = len(kpis) if len(kpis) > 0 else 1
            for i, kpi in enumerate(kpis):
                left = Inches(1 + (i * (11.33/num_kpis)))
                box = slide.shapes.add_shape(1, left, Inches(2.5), Inches(3.5), Inches(2.5))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(245, 245, 247) if item.get('theme') != 'dark' else RGBColor(39, 39, 41)
                tf = box.text_frame
                tf.text = kpi.get('value', '')
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(60)
                p.font.bold = True
                p.font.color.rgb = RGBColor(3, 78, 162) # Samsung Blue
                p2 = tf.add_paragraph()
                p2.text = kpi.get('label', '')
                p2.alignment = PP_ALIGN.CENTER
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(29, 29, 31) if item.get('theme') != 'dark' else RGBColor(255, 255, 255)

        elif stype == 'table':
            rows_data = item.get('rows', [])
            headers = item.get('headers', [])
            rows, cols = len(rows_data) + 1, len(headers)
            table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.5), Inches(11.33), Inches(4)).table
            for c, h in enumerate(headers): table.cell(0, c).text = h
            for r, row in enumerate(rows_data):
                for c, val in enumerate(row): table.cell(r+1, c).text = val

        elif stype == 'flow':
            steps = item.get('steps', [])
            for i, step in enumerate(steps):
                left = Inches(1 + (i * 3.5))
                box = slide.shapes.add_shape(1, left, Inches(3.5), Inches(2.5), Inches(1.5))
                box.text = step
                if i < len(steps) - 1:
                    slide.shapes.add_textbox(left + Inches(2.6), Inches(4), Inches(0.8), Inches(0.5)).text = "→"

        elif stype == 'comparison':
            comps = item.get('comparisons', [])
            for i, comp in enumerate(comps):
                left = Inches(1 + (i * 6))
                box = slide.shapes.add_shape(1, left, Inches(2.5), Inches(5), Inches(4))
                box.text = f"{comp.get('label')}\n\n{comp.get('text')}"
        
        elif stype == 'quadrant':
            quads = item.get('quadrants', [])
            for i, q in enumerate(quads):
                left = Inches(2 + (i%2 * 5))
                top = Inches(2.5 + (i//2 * 2.2))
                box = slide.shapes.add_shape(1, left, top, Inches(4.5), Inches(2))
                box.text = f"{q.get('label')}\n{q.get('text')}"
                
        elif stype == 'swot':
            swot = item.get('swot', {})
            keys = ['S', 'W', 'O', 'T']
            for i, k in enumerate(keys):
                left = Inches(1.5 + (i%2 * 5.5))
                top = Inches(2.5 + (i//2 * 2.2))
                box = slide.shapes.add_shape(1, left, top, Inches(5), Inches(2))
                box.text = f"{k}\n{swot.get(k, '')}"

        else:
            body = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(4))
            body.text_frame.text = item.get('content', '')

    try:
        prs.save(output_pptx)
        print(f"High-Fidelity PPT Export Successful: {output_pptx}")
    except Exception as e:
        print(f"Export Failed: {e}")

if __name__ == "__main__":
    json_path = sys.argv[1] if len(sys.argv) > 1 else "slides_data.json"
    export_high_fidelity_ppt(json_path)
